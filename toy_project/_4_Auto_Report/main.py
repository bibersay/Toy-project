from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)

title = slide.placeholders[0]
title.text = "hello world"

subtitle = slide.placeholders[1]
subtitle.text = "python-pptx was here"

prs.save('test.pptx')



bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)

title_shape = slide.placeholders[0]
title_shape.text = 'adding a bullet slide'

body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = 'find ghe bullet slide layout'

p = tf.add_paragraph()
p.text = 'use_textframe.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'use _textframe.add_paragraph() for subsequent bullets'
p.level = 2

prs.save('test.pptx')